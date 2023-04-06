import pytest
from pptx import Presentation
from pptx.shapes.picture import Picture
from pptx.enum.shapes import MSO_SHAPE_TYPE
import glob
from pptx.table import Table
from pptx.util import Inches

def iter_picture_shapes(prs):
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                yield shape

                
                
def test_iter_picture_shapes():
    # create a presentation object with at least one slide containing a picture shape
    prs = Presentation('presentation.pptx')
    
    # call the iter_picture_shapes function on the presentation object
    picture_shapes = list(iter_picture_shapes(prs))
    
    # assert that the picture_shapes object is a list
    assert isinstance(picture_shapes, list)
    
    # assert that all elements in the picture_shapes list are Picture objects
    assert all(isinstance(shape, Picture) for shape in picture_shapes)
    
    # assert that the picture_shapes list is not empty
    assert len(picture_shapes) > 0
    
def iter_to_nonempty_table_cells(tbl):
    for ridx in range(sum(1 for _ in iter(tbl.rows))):
        for cidx in range(sum(1 for _ in iter(tbl.columns))):
            cell = tbl.cell(ridx, cidx)
            txt = type("")(cell.text)
            txt = txt.strip()
            yield txt  


def get_tables_from_presentation(pres):
    tables = list()
    for slide in pres.slides:
        for shp in iter(slide.shapes):
            if shp.has_table:
                table = shp.table
                tables.append(table)
    return tables


def test_table_data_extraction():
    # create a presentation with a slide containing a table
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    table = slide.shapes.add_table(rows=2, cols=2, left=Inches(1), top=Inches(2), width=Inches(4), height=Inches(1)).table
    table.cell(0, 0).text = "Name"
    table.cell(0, 1).text = "Age"
    table.cell(1, 0).text = "John"
    table.cell(1, 1).text = "30"
    
    # extract table data
    table_data = []
    tables = get_tables_from_presentation(prs)
    for tbl in tables:
        it = iter_to_nonempty_table_cells(tbl)
        table_rows = []
        for i in range(sum(1 for _ in iter(tbl.rows))):
            row_data = []
            for j in range(sum(1 for _ in iter(tbl.columns))):
                row_data.append(next(it))
            table_rows.append(row_data)
        table_data.append(table_rows)
    
    # check if the extracted data matches the expected data
    assert table_data == [[["Name", "Age"], ["John", "30"]]]

    
class TestTextDataExtraction(unittest.TestCase):
    def setUp(self):
        self.presentation = Presentation('presentation.pptx')
        
    def test_text_extraction(self):
        slide_count = 0
        for slide in self.presentation.slides:
            print("----------------------")
            print("S: ", slide_count)
            slide_count += 1

            #Find the heading paragraph in the slide
            heading=None
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    if shape.text.startswith("#"):
                        heading=shape.text.strip("#")
                        break
            if heading is not None:
                print("Heading:",heading)

            #find the sub points for the heading

                sub_points=[]
                for shape in slide.shapes:
                    if hasattr(shape,"text"):
                        if shape.text.startswith("-"):
                            sub_points.append(shape.textstrip("-"))

                if len(sub_points)>0:
                    print("Sub Points:")
                    for sub_point in sub_points:
                        print("- " +sub_point)

            #print the text in shapes
            for shape in slide.shapes:
                if hasattr(shape,"text"):
                    if not shape.text.startswith("#")and not shape.text.startswith("-"):
                        print(shape.text)
                        print()

if __name__ == '__main__':
    unittest.main()    
